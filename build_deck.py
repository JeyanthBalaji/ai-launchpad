"""Build the AI Launchpad pitch deck (Track 3, Unicorn) with python-pptx.
Slide 5 shows the real startup kit generated on an AMD Developer Cloud GPU.
Saves AI_Launchpad_Deck.pptx next to this script."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BG=RGBColor(0x0F,0x17,0x2A); BG2=RGBColor(0x0B,0x11,0x20)
CARD=RGBColor(0x1E,0x29,0x3B); CARD2=RGBColor(0x24,0x31,0x47)
SKY=RGBColor(0x38,0xBD,0xF8); BLUE=RGBColor(0x25,0x63,0xEB); INDIGO=RGBColor(0x63,0x66,0xF1)
WHITE=RGBColor(0xF8,0xFA,0xFC); MUTED=RGBColor(0x94,0xA3,0xB8); EMER=RGBColor(0x34,0xD3,0x99)
AMBER=RGBColor(0xF5,0x9E,0x0B); LINE=RGBColor(0x33,0x41,0x55)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=13.333,7.5

def slide(bg=BG):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s

def txt(s,x,y,w,h,runs,size=16,color=WHITE,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,italic=False,font="Calibri",spacing=1.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    if isinstance(runs,str): runs=[(runs,color,bold)]
    p=tf.paragraphs[0]; p.alignment=align
    if spacing!=1.0: p.line_spacing=spacing
    for t,c,b in runs:
        r=p.add_run(); r.text=t; r.font.size=Pt(size); r.font.color.rgb=c; r.font.bold=b
        r.font.italic=italic; r.font.name=font
    return tb

def card(s,x,y,w,h,fill=CARD,radius=0.09,shadow=True):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=LINE; sh.line.width=Pt(0.75)
    try: sh.adjustments[0]=radius
    except Exception: pass
    sh.shadow.inherit=False
    if shadow:
        el=sh._element.spPr; ef=el.makeelement(qn('a:effectLst'),{})
        sdw=ef.makeelement(qn('a:outerShdw'),{'blurRad':'90000','dist':'40000','dir':'5400000','rotWithShape':'0'})
        clr=sdw.makeelement(qn('a:srgbClr'),{'val':'000000'}); alpha=clr.makeelement(qn('a:alpha'),{'val':'32000'})
        clr.append(alpha); sdw.append(clr); ef.append(sdw); el.append(ef)
    return sh

def circle(s,x,y,d,fill,ch,chcolor=WHITE,chsize=18):
    c=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb=fill; c.line.fill.background(); c.shadow.inherit=False
    tf=c.text_frame; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=ch; r.font.size=Pt(chsize); r.font.bold=True; r.font.color.rgb=chcolor; r.font.name="Calibri"
    return c

def rrect(s,x,y,w,h,fill,text=None,tcolor=WHITE,tsize=12,bold=True,radius=0.5,line=None):
    p=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    p.fill.solid(); p.fill.fore_color.rgb=fill
    if line is None: p.line.fill.background()
    else: p.line.color.rgb=line; p.line.width=Pt(0.75)
    try: p.adjustments[0]=radius
    except Exception: pass
    p.shadow.inherit=False
    if text is not None:
        tf=p.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.05); tf.margin_right=Inches(0.05); tf.margin_top=0; tf.margin_bottom=0
        tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        par=tf.paragraphs[0]; par.alignment=PP_ALIGN.CENTER
        r=par.add_run(); r.text=text; r.font.size=Pt(tsize); r.font.color.rgb=tcolor; r.font.name="Calibri"; r.font.bold=bold
    return p

def pill(s,x,y,w,text):
    p=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(0.42))
    p.fill.solid(); p.fill.fore_color.rgb=CARD; p.line.color.rgb=LINE; p.line.width=Pt(0.75)
    try: p.adjustments[0]=0.5
    except Exception: pass
    p.shadow.inherit=False
    tf=p.text_frame; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    par=tf.paragraphs[0]; par.alignment=PP_ALIGN.CENTER
    r=par.add_run(); r.text="●  "+text; r.font.size=Pt(11); r.font.color.rgb=SKY; r.font.name="Calibri"; r.font.bold=True
    return p

# S1 TITLE
s=slide()
pill(s,SW/2-2.15,1.35,4.3,"Powered by AMD Developer Cloud GPUs")
txt(s,1,2.35,11.333,1.3,[("AI ",WHITE,True),("Launchpad",SKY,True)],size=60,align=PP_ALIGN.CENTER)
txt(s,1.5,3.75,10.333,0.7,"Turn one idea into a full startup kit in seconds.",size=22,color=MUTED,align=PP_ALIGN.CENTER)
txt(s,1,5.55,11.333,0.4,"AMD Developer Hackathon: ACT II  ·  Unicorn Track (Open Innovation)",size=15,color=SKY,align=PP_ALIGN.CENTER,bold=True)
txt(s,1,6.05,11.333,0.4,"Jeyanth Balaji K",size=14,color=MUTED,align=PP_ALIGN.CENTER)

# S2 PROBLEM
s=slide()
txt(s,0.9,0.7,11.5,0.8,"The blank-page problem",size=38,bold=True)
txt(s,0.9,1.55,11.5,0.7,"Starting a business means hours of work before you can show anyone your idea.",size=18,color=MUTED)
probs=[("1","Naming","Brainstorming a name, tagline, and value proposition that actually lands.",BLUE),
("2","Copywriting","Writing a landing page — headline, features, and a call to action.",INDIGO),
("3","Branding","Choosing colors, a logo direction, and a consistent visual identity.",SKY)]
cx,cw,gap=0.9,3.74,0.30
for i,(n,title,desc,col) in enumerate(probs):
    x=cx+i*(cw+gap); card(s,x,2.65,cw,2.6); circle(s,x+0.35,3.0,0.7,col,n,WHITE,22)
    txt(s,x+0.35,3.95,cw-0.7,0.5,title,size=20,bold=True)
    txt(s,x+0.35,4.5,cw-0.7,1.3,desc,size=14,color=MUTED,spacing=1.05)
txt(s,0.9,5.7,11.5,0.9,[("The result: ",MUTED,False),("great ideas stall at the starting line.",WHITE,True)],size=20,align=PP_ALIGN.CENTER)

# S3 SOLUTION
s=slide()
txt(s,0.9,0.7,11.5,0.8,"One sentence in. A full kit out.",size=38,bold=True)
txt(s,0.9,1.55,11.5,0.9,[("Type one line — ",MUTED,False),("“an app that helps students find part-time jobs”",WHITE,True),(" — and a crew of AI agents delivers a complete starter kit in ~30 seconds.",MUTED,False)],size=18,spacing=1.1)
outs=[("Name & tagline","A punchy identity + value proposition",BLUE),("Landing page","Hero, three features, and a call to action",INDIGO),("Brand look","A 5-color palette, logo concept, and font",SKY),("Social captions","Three launch posts, ready to publish",EMER)]
cw,gap=2.83,0.28; x0=0.9
for i,(t,d,col) in enumerate(outs):
    x=x0+i*(cw+gap); card(s,x,2.9,cw,2.5)
    bar=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.32),Inches(3.2),Inches(0.42),Inches(0.42))
    bar.fill.solid(); bar.fill.fore_color.rgb=col; bar.line.fill.background(); bar.shadow.inherit=False
    txt(s,x+0.32,3.85,cw-0.6,0.55,t,size=17,bold=True); txt(s,x+0.32,4.4,cw-0.6,1.0,d,size=13.5,color=MUTED,spacing=1.05)
txt(s,0.9,5.95,11.5,0.7,[("You watch it appear live — ",MUTED,False),("the “wow” moment.",SKY,True)],size=19,align=PP_ALIGN.CENTER)

# S4 HOW IT WORKS
s=slide()
txt(s,0.9,0.7,11.5,0.8,"A crew of four AI agents",size=38,bold=True)
txt(s,0.9,1.55,11.5,0.6,"The Strategist runs first — then the other three work in parallel for speed.",size=17,color=MUTED)
agents=[("Strategist","name · tagline · value prop",BLUE),("Copywriter","landing-page copy",INDIGO),("Brand Designer","palette · logo · font",SKY),("Social Manager","3 launch captions",EMER)]
cw,gap=2.83,0.28; x0=0.9
for i,(t,d,col) in enumerate(agents):
    x=x0+i*(cw+gap); card(s,x,2.5,cw,2.15); circle(s,x+0.32,2.82,0.62,col,str(i+1),WHITE,20)
    txt(s,x+1.05,2.86,cw-1.2,0.6,t,size=16,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.32,3.65,cw-0.6,0.8,d,size=13,color=MUTED,spacing=1.05)
card(s,0.9,5.05,11.533,1.55,fill=CARD2)
txt(s,1.25,5.28,10.9,1.15,[("Under the hood:  ",SKY,True),("each agent is a single focused LLM call on ",MUTED,False),("an open model running on AMD GPUs (vLLM)",WHITE,True),(", forced to return JSON so the UI renders reliably. The orchestrator is ",MUTED,False),("resilient",WHITE,True),(" — if one agent fails, the others still return, so a live demo never breaks.",MUTED,False)],size=15.5,anchor=MSO_ANCHOR.MIDDLE,spacing=1.15)

# S5 REAL OUTPUT (generated on AMD)
s=slide()
txt(s,0.9,0.62,11.5,0.8,"Real output, generated on AMD",size=38,bold=True)
txt(s,0.9,1.45,11.5,0.55,[("From ",MUTED,False),("“an app that helps students find part-time jobs”",WHITE,True),(" — one real run of the crew on an AMD GPU.",MUTED,False)],size=16,spacing=1.05)

# Left: landing-page mock
card(s,0.9,2.15,7.4,4.35,fill=CARD,radius=0.04)
rrect(s,1.25,2.5,0.72,0.72,BLUE,"J",WHITE,26,True,0.25)
txt(s,2.12,2.5,5.2,0.5,"JobHopper",size=23,bold=True,anchor=MSO_ANCHOR.MIDDLE)
txt(s,2.14,3.04,5.6,0.35,"Find your flex job, college student!",size=12.5,color=SKY,italic=True)
txt(s,1.25,3.62,6.85,0.7,"Find Your Dream Job, Tailored Just For You",size=18,bold=True,spacing=1.02)
txt(s,1.25,4.42,6.85,0.75,"JobHopper connects you with customized job opportunities that fit your skills and preferences. Work on your terms.",size=12,color=MUTED,spacing=1.05)
chips=["Tailored Opportunities","Flexible Hours","Personalized Recs"]
cwc=2.18
for i,c in enumerate(chips):
    rrect(s,1.25+i*(cwc+0.15),5.28,cwc,0.5,CARD2,c,WHITE,10.5,True,0.5,line=LINE)
rrect(s,1.25,5.95,4.25,0.5,BLUE,"Start Your JobHopper Journey Now  →",WHITE,12,True,0.5)

# Right: brand + social + engine
card(s,8.55,2.15,3.88,4.35,fill=CARD2)
txt(s,8.9,2.42,3.3,0.35,"BRAND PALETTE",size=11,bold=True,color=SKY)
palette=[RGBColor(0x4C,0xAF,0x50),RGBColor(0x2E,0x8B,0x57),RGBColor(0x32,0xCD,0x32),RGBColor(0x98,0xFB,0x98),RGBColor(0xD3,0xD3,0xD3)]
for i,col in enumerate(palette):
    sw=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(8.9+i*0.58),Inches(2.8),Inches(0.5),Inches(0.5))
    sw.fill.solid(); sw.fill.fore_color.rgb=col; sw.line.color.rgb=LINE; sw.line.width=Pt(0.5); sw.shadow.inherit=False
    try: sw.adjustments[0]=0.18
    except Exception: pass
txt(s,8.9,3.55,3.3,0.35,"LAUNCH POST",size=11,bold=True,color=SKY)
txt(s,8.9,3.9,3.3,1.5,"“Revolutionize your job search with JobHopper! Find tailored opportunities & work on your terms. #JobHopper #CareerRevolution”",size=11,color=WHITE,spacing=1.08,italic=True)
rrect(s,8.9,5.85,3.18,0.5,BG,"▸ Qwen2.5-7B-Instruct · vLLM on AMD GPU",SKY,10,True,0.5,line=LINE)

txt(s,0.9,6.66,11.5,0.5,[("✓  ",EMER,True),("Generated live on an AMD Developer Cloud GPU — no cloud LLM API, no mock data.",MUTED,False)],size=12.5,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# S6 BUILT ON AMD
s=slide()
txt(s,0.9,0.7,11.5,0.8,"Built on AMD compute",size=38,bold=True)
txt(s,0.9,1.55,11.5,0.6,"Every agent runs on AMD GPUs — and the whole app is containerized for AMD Developer Cloud.",size=17,color=MUTED)
stack=[("Frontend","React + Vite + Tailwind CSS",BLUE),("Backend","Python + FastAPI orchestration",INDIGO),("AI compute","AMD Developer Cloud GPU · vLLM",SKY),("Packaging","Docker → AMD Developer Cloud",EMER)]
cw,gap=2.83,0.28; x0=0.9
for i,(t,d,col) in enumerate(stack):
    x=x0+i*(cw+gap); card(s,x,2.6,cw,2.35)
    bar=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+0.32),Inches(2.92),Inches(0.42),Inches(0.42))
    bar.fill.solid(); bar.fill.fore_color.rgb=col; bar.line.fill.background(); bar.shadow.inherit=False
    txt(s,x+0.32,3.55,cw-0.6,0.5,t,size=17,bold=True); txt(s,x+0.32,4.08,cw-0.6,0.9,d,size=13.5,color=MUTED,spacing=1.05)
card(s,0.9,5.35,11.533,1.2,fill=CARD2)
txt(s,1.25,5.5,10.9,0.9,[("Proven on an AMD Developer Cloud GPU (Qwen2.5-7B via vLLM) · ",MUTED,False),("same OpenAI-compatible code runs on Fireworks AI, also on AMD GPUs.",WHITE,True)],size=15.5,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER,spacing=1.05)

# S7 WHY IT STANDS OUT
s=slide()
txt(s,0.9,0.65,11.5,0.8,"Why AI Launchpad stands out",size=38,bold=True)
feats=[("Live agent reveal","Judges watch the crew work, one by one.",BLUE),("Full brand kit","Name, page, palette, logo mark, captions.",INDIGO),("Download kit","Exports a standalone, brand-styled landing page.",SKY),("Works for any idea","Coffee, fitness, tutoring — unique every time.",EMER),("Resilient by design","One failed agent never breaks the demo.",AMBER),("Runs on AMD GPUs","Self-hosted vLLM on AMD Developer Cloud.",BLUE)]
cw,ch,gx,gy=3.74,1.75,0.30,0.30; x0,y0=0.9,1.7
for i,(t,d,col) in enumerate(feats):
    r,c=divmod(i,3); x=x0+c*(cw+gx); y=y0+r*(ch+gy); card(s,x,y,cw,ch)
    circle(s,x+0.3,y+0.32,0.5,col,"✓",WHITE,16)
    txt(s,x+0.98,y+0.32,cw-1.15,0.5,t,size=15.5,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.3,y+0.98,cw-0.55,0.7,d,size=12.5,color=MUTED,spacing=1.05)

# S8 VISION
s=slide(BG2)
txt(s,1,2.35,11.333,1.1,"The fastest way from idea to launch.",size=40,bold=True,align=PP_ALIGN.CENTER)
txt(s,1,3.6,11.333,0.6,"For every founder, student, and creator.",size=20,color=SKY,align=PP_ALIGN.CENTER)
pill(s,SW/2-2.6,4.75,5.2,"github.com/JeyanthBalaji/ai-launchpad")
txt(s,1,5.75,11.333,0.6,[("AI ",WHITE,True),("Launchpad",SKY,True),("  ·  AMD Developer Hackathon: ACT II",MUTED,False)],size=16,align=PP_ALIGN.CENTER,bold=True)

outdir=os.path.dirname(os.path.abspath(__file__))
prs.save(os.path.join(outdir,"AI_Launchpad_Deck.pptx"))
print("saved", len(prs.slides._sldIdLst), "slides to", outdir)
